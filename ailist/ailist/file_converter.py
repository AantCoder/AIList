"""
file_converter.py -- File Converter MCP Server
==============================================

Локальный MCP-сервер для конвертации файлов в текст и экспорта документов.
Работает с файлами на локальной файловой системе.

Поддерживаемые форматы (чтение -> текст)
-----------------------------------------
  PDF                     -> текст (через pypdf)
  Word  (.docx, .doc)     -> текст (через python-docx / LibreOffice)
  Excel (.xlsx, .xls)     -> текст (через openpyxl)
  PowerPoint (.pptx)      -> текст (через python-pptx)
  HTML  (.html, .htm)     -> текст (через BeautifulSoup)
  Аудио (.mp3, .wav, .ogg, .m4a, .flac, .aac) -> транскрипция (через faster-whisper, локально)
  Видео (.mp4, .mov, .avi, .webm) -> транскрипция аудиодорожки + кадры в base64
  Изображения (.png, .jpg, .jpeg, .webp, .gif, .bmp) -> текстовое описание (через LLaVA в Ollama)

Поддерживаемые форматы (экспорт -> файл)
-----------------------------------------
  Markdown / HTML / текст -> PDF      (через xhtml2pdf)
  Структурированный текст -> DOCX     (через python-docx)
  Список строк/словарей   -> XLSX     (через openpyxl)

Установка и запуск как MCP-сервера
-----------------------------------

  # 1. Установить зависимости (только нужные форматы):
  pip install fastmcp
  pip install pypdf              # PDF -> текст
  pip install python-docx        # .docx / .doc (для .doc ещё нужен LibreOffice)
  pip install openpyxl           # .xlsx / .xls
  pip install python-pptx        # .pptx
  pip install "beautifulsoup4[lxml]"  # HTML
  pip install faster-whisper     # аудио и видео -> транскрипция
  pip install opencv-python      # видео -> кадры
  pip install moviepy            # видео -> аудио -> транскрипция
  pip install xhtml2pdf         # текст/markdown/HTML -> PDF
  # изображения -> описание через LLaVA:
  #   ollama pull llava           # скачать модель (~4 ГБ), нужен запущенный Ollama

  # 2. Запустить сервер:
  python file_converter.py                        # SSE на 127.0.0.1:8010
  python file_converter.py --port 9000            # другой порт
  python file_converter.py --host 0.0.0.0        # доступен в сети
  python file_converter.py --cuda                 # GPU для Whisper
  python file_converter.py --model small          # другой размер модели Whisper

  # Либо через fastmcp CLI (stdio-транспорт, для Claude Desktop):
  fastmcp run file_converter.py

Подключение из разных клиентов
--------------------------------

  # Claude Desktop (claude_desktop_config.json):
  {
    "mcpServers": {
      "file-converter": {
        "command": "python",
        "args": ["/path/to/file_converter.py"]
      }
    }
  }

  # Claude Desktop с аргументами (GPU, другая модель):
  {
    "mcpServers": {
      "file-converter": {
        "command": "python",
        "args": ["/path/to/file_converter.py", "--cuda", "--model", "large"]
      }
    }
  }

  # Claude Desktop через переменные окружения:
  {
    "mcpServers": {
      "file-converter": {
        "command": "python",
        "args": ["/path/to/file_converter.py"],
        "env": {
          "FILE_CONVERTER_CUDA": "1",
          "FILE_CONVERTER_MODEL": "large",
          "FILE_CONVERTER_LANGUAGE": "ru",
          "FILE_CONVERTER_LLAVA_MODEL": "llava:7b",
          "FILE_CONVERTER_OLLAMA_URL": "http://localhost:11434"
        }
      }
    }
  }

  # LangChain / AIList через SSE (после запуска сервера):
  await ai.mcp_connect("file-converter", url="http://127.0.0.1:8010/sse")

  # Любой MCP-клиент через stdio (без запуска сервера вручную):
  command: python file_converter.py
  transport: stdio

Инструменты сервера
--------------------
  convert_pdf(path)                                      -> str
  convert_office(path)                                   -> str
  convert_html(path)                                     -> str
  transcribe_audio(path, language?, model_size?, word_timestamps?) -> str
  transcribe_video(path, language?)                      -> str
  extract_video_frames(path, num_frames?)                -> str  # JSON-массив base64-JPEG
  describe_image(path, prompt?, model?)                  -> str  # требует Ollama + LLaVA
  export_pdf(content, output_path?, title?)              -> str  # markdown/текст -> PDF
  export_docx(content, output_path?, title?)             -> str  # текст -> DOCX
  export_xlsx(rows, output_path?, sheet_name?, headers?) -> str  # данные -> XLSX

  Все инструменты принимают абсолютный путь к файлу на той машине, где запущен сервер.
"""

from __future__ import annotations

import base64
import os
import tempfile
from pathlib import Path

# -- Опциональные библиотеки конвертации --------------------------------------
# Каждая подгружается только если реально используется соответствующий формат.
# Отсутствие библиотеки не мешает запуску сервера -- только вызов нужного
# инструмента завершится ошибкой с понятным сообщением об установке.

try:
    from pypdf import PdfReader as _PdfReader
except ImportError:
    _PdfReader = None

try:
    import docx as _docx
except ImportError:
    _docx = None

try:
    import openpyxl as _openpyxl
except ImportError:
    _openpyxl = None

try:
    from pptx import Presentation as _Presentation
except ImportError:
    _Presentation = None

try:
    from bs4 import BeautifulSoup as _BeautifulSoup
except ImportError:
    _BeautifulSoup = None

try:
    from faster_whisper import WhisperModel as _WhisperModel
except ImportError:
    _WhisperModel = None

try:
    import cv2 as _cv2
except ImportError:
    _cv2 = None

try:
    import xhtml2pdf.pisa as _pisa
except ImportError:
    _pisa = None

try:
    from moviepy.editor import VideoFileClip as _VideoFileClip  # type: ignore
except ImportError:
    try:
        from moviepy import VideoFileClip as _VideoFileClip
    except ImportError:
        _VideoFileClip = None

DEFAULT_MODEL = "llava:7b"  # по умолчанию используем более лёгкую модель LLaVA, которая работает на CPU и занимает ~4 ГБ, вместо llava:13b (~8 ГБ)

# -- Сообщения об ошибках ------------------------------------------------------

_ERR_PDF_NO_LIB             = "install pypdf to read '{name}': pip install pypdf"
_ERR_OFFICE_NO_LIB          = "install required library to read '{name}': pip install python-docx openpyxl python-pptx"
_ERR_HTML_NO_LIB            = "install beautifulsoup4 to read '{name}': pip install \"beautifulsoup4[lxml]\""
_ERR_AUDIO_NO_LIB           = "install faster-whisper to transcribe '{name}': pip install faster-whisper"
_ERR_VIDEO_NO_CV2           = "install opencv-python to extract frames from '{name}': pip install opencv-python"
_ERR_VIDEO_CANT_OPEN        = "cannot open video '{name}' -- file is corrupted or format is not supported"
_ERR_VIDEO_NO_MOVIEPY       = "install moviepy to extract audio from '{name}': pip install moviepy"
_ERR_DOC_NO_LIBREOFFICE     = "install LibreOffice to read .doc files, or convert '{name}' to .docx first"
_ERR_DOC_LIBREOFFICE_FAILED = "LibreOffice failed to convert '{name}' to .docx -- file may be corrupted"
_ERR_LLAVA_OLLAMA_UNAVAIL   = "Ollama is not running or LLaVA model is not available. Start Ollama and run: ollama pull {model}"

_ERR_EXPORT_PDF_NO_LIB  = "install xhtml2pdf to export PDF: pip install xhtml2pdf"
_ERR_EXPORT_DOCX_NO_LIB = "install python-docx to export DOCX: pip install python-docx"
_ERR_EXPORT_XLSX_NO_LIB = "install openpyxl to export XLSX: pip install openpyxl"

_OLLAMA_KEEPALIVE = "5m"   # как долго Ollama держит модель загруженной после вызова

_AUDIO_WORD_TIMESTAMP = "[{start:.2f}s -> {end:.2f}s] {word}"
_OFFICE_SHEET_LABEL   = "[Sheet: {title}]"
_OFFICE_SLIDE_LABEL   = "[Slide {num}]"


# =============================================================================
# Класс-конвертер -- чистая логика без зависимости от MCP или LangChain
# =============================================================================

class FileConverter:
    """
    Конвертирует файлы разных форматов в текст или base64.

    Используется напрямую или через MCP-обёртку (см. модульный объект `mcp` ниже).

    Параметры:
        use_cuda         -- использовать GPU для Whisper (требует CUDA, ~3-4 ГБ VRAM).
        whisper_model    -- размер модели Whisper: tiny/base/small/medium/large/turbo.
        whisper_language -- язык по умолчанию (None = автодетект).
        llava_model      -- модель LLaVA в Ollama (по умолчанию "llava:7b").
                           Другие vision-модели Ollama тоже работают: llava:13b, llava:34b,
                           llava-llama3, bakllava и т.д.
        ollama_url       -- адрес Ollama API (по умолчанию http://localhost:11434).
        llava_cpu_only   -- принудительно запускать vision-модель на CPU (num_gpu=0).
                           По умолчанию True: GPU не используется, чтобы не вытеснять
                           основную LLM из видеопамяти. Каждый вызов передаёт
                           keep_alive=_OLLAMA_KEEPALIVE -- Ollama держит модель
                           загруженной между вызовами без повторной инициализации.
    """

    def __init__(
        self,
        use_cuda: bool = False,
        whisper_model: str = "turbo",
        whisper_language: str | None = None,
        llava_model: str = DEFAULT_MODEL,
        ollama_url: str = "http://localhost:11434",
        llava_cpu_only: bool = True,
        cache_dir: "Path | None" = None,
    ):
        self.use_cuda = use_cuda
        self.whisper_model = whisper_model
        self.whisper_language = whisper_language
        self.llava_model = llava_model
        self.ollama_url = ollama_url.rstrip("/")
        self.llava_cpu_only = llava_cpu_only
        self.cache_dir = cache_dir  # папка для загрузки Whisper-моделей; None = дефолт faster-whisper (~/.cache/huggingface)
        self._whisper_cache: dict = {}  # (model_size, device) -> WhisperModel

    def pdf_to_text(self, path: str) -> str:
        """Извлекает текст из PDF постранично через pypdf."""
        if _PdfReader is None:
            raise RuntimeError(_ERR_PDF_NO_LIB.format(name=Path(path).name))
        reader = _PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(p for p in pages if p.strip())

    def office_to_text(self, path: str, ext: str | None = None) -> str:
        """
        Извлекает текст из Office-документа.
        Формат определяется по расширению файла (или явно через ext).
        Поддерживает: .docx, .doc, .xlsx, .xls, .pptx.
        """
        path = str(path)
        name = Path(path).name
        ext  = (ext or Path(path).suffix.lstrip(".")).lower()

        if ext == "docx":
            if _docx is None:
                raise RuntimeError(_ERR_OFFICE_NO_LIB.format(name=name))
            doc = _docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        if ext == "doc":
            # Старый бинарный формат -- конвертируем через LibreOffice
            import subprocess, shutil
            if shutil.which("libreoffice") is None:
                raise RuntimeError(_ERR_DOC_NO_LIBREOFFICE.format(name=name))
            with tempfile.TemporaryDirectory() as tmp:
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", tmp, path],
                    check=True, capture_output=True,
                )
                converted = list(Path(tmp).glob("*.docx"))
                if not converted:
                    raise RuntimeError(_ERR_DOC_LIBREOFFICE_FAILED.format(name=name))
                if _docx is None:
                    raise RuntimeError(_ERR_OFFICE_NO_LIB.format(name=name))
                doc = _docx.Document(str(converted[0]))
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        if ext in ("xlsx", "xls"):
            if _openpyxl is None:
                raise RuntimeError(_ERR_OFFICE_NO_LIB.format(name=name))
            wb = _openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(_OFFICE_SHEET_LABEL.format(title=sheet.title))
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        parts.append("\t".join(cells))
            wb.close()
            return "\n".join(parts)

        if ext == "pptx":
            if _Presentation is None:
                raise RuntimeError(_ERR_OFFICE_NO_LIB.format(name=name))
            prs = _Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(_OFFICE_SLIDE_LABEL.format(num=i))
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)

        raise RuntimeError(_ERR_OFFICE_NO_LIB.format(name=name))

    def html_to_text(self, path: str) -> str:
        """Извлекает читаемый текст из HTML, убирая теги, скрипты и стили."""
        if _BeautifulSoup is None:
            raise RuntimeError(_ERR_HTML_NO_LIB.format(name=Path(path).name))
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read()
        soup = _BeautifulSoup(raw, "lxml")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)

    def audio_to_transcript(
        self,
        path: str,
        *,
        model_size: str | None = None,
        language: str | None = None,
        word_timestamps: bool = False,
        vad_filter: bool = True,
        initial_prompt: str | None = None,
    ) -> str:
        """
        Транскрибирует аудио локально через faster-whisper (без API, без интернета).

        Модель скачивается с HuggingFace автоматически при первом вызове.
        Последующие вызовы используют кэшированную модель -- без скачивания.

        Параметры:
            model_size      -- tiny/base/small/medium/large/turbo (умолчание: turbo).
            language        -- ISO 639-1 ("ru", "en", ...). None = автодетект.
            word_timestamps -- True -> строки вида "[0.00s -> 0.54s] слово".
            vad_filter      -- фильтровать тишину через Silero VAD (рекомендуется).
            initial_prompt  -- подсказка с терминологией или контекстом.
        """
        name = Path(path).name
        if _WhisperModel is None:
            raise RuntimeError(_ERR_AUDIO_NO_LIB.format(name=name))

        model_size = model_size or self.whisper_model
        language   = language   or self.whisper_language
        device     = "cuda" if self.use_cuda else "cpu"
        cache_key  = (model_size, device)

        if cache_key not in self._whisper_cache:
            compute_type = "float16" if device == "cuda" else "int8"
            kwargs = dict(device=device, compute_type=compute_type)
            if self.cache_dir is not None:
                self.cache_dir.mkdir(parents=True, exist_ok=True)
                kwargs["download_root"] = str(self.cache_dir)
            self._whisper_cache[cache_key] = _WhisperModel(model_size, **kwargs)

        segments, _ = self._whisper_cache[cache_key].transcribe(
            str(path),
            language=language,
            word_timestamps=word_timestamps,
            vad_filter=vad_filter,
            initial_prompt=initial_prompt,
        )

        if word_timestamps:
            lines = []
            for seg in segments:
                for word in seg.words:
                    lines.append(_AUDIO_WORD_TIMESTAMP.format(
                        start=word.start, end=word.end, word=word.word,
                    ))
            return "\n".join(lines)
        return " ".join(seg.text.strip() for seg in segments)

    def video_extract_frames_b64(self, path: str, num_frames: int = 5) -> list[str]:
        """
        Извлекает num_frames равномерно распределённых кадров из видео.
        Возвращает список строк в формате base64-JPEG.
        """
        name = Path(path).name
        if _cv2 is None:
            raise RuntimeError(_ERR_VIDEO_NO_CV2.format(name=name))
        cap = _cv2.VideoCapture(str(path))
        if not cap.isOpened():
            raise RuntimeError(_ERR_VIDEO_CANT_OPEN.format(name=name))
        total = int(cap.get(_cv2.CAP_PROP_FRAME_COUNT))
        if total <= 0:
            cap.release()
            raise RuntimeError(_ERR_VIDEO_CANT_OPEN.format(name=name))
        step = max(1, total // num_frames)
        frames = []
        for i in range(num_frames):
            cap.set(_cv2.CAP_PROP_POS_FRAMES, i * step)
            ret, frame = cap.read()
            if not ret:
                break
            _, buf = _cv2.imencode(".jpg", frame)
            frames.append(base64.b64encode(buf).decode("utf-8"))
        cap.release()
        return frames

    def video_extract_audio_transcript(self, path: str) -> str | None:
        """
        Извлекает аудиодорожку из видео и транскрибирует через Whisper.

        Возвращает:
            str  -- транскрипт (может быть пустым если речи нет).
            None -- аудиодорожка отсутствует (не ошибка).
        """
        name = Path(path).name
        if _VideoFileClip is None:
            raise RuntimeError(_ERR_VIDEO_NO_MOVIEPY.format(name=name))
        if _WhisperModel is None:
            raise RuntimeError(_ERR_AUDIO_NO_LIB.format(name=name))

        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
                tmp_path = tmp.name
            clip = _VideoFileClip(str(path))
            if clip.audio is None:
                clip.close()
                return None
            try:
                clip.audio.write_audiofile(tmp_path, logger=None)
            except TypeError:
                clip.audio.write_audiofile(tmp_path)
            clip.close()
            return self.audio_to_transcript(tmp_path)
        except RuntimeError:
            raise
        except Exception as e:
            raise RuntimeError(_ERR_AUDIO_NO_LIB.format(name=name)) from e
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

    def describe_image(
        self,
        path: str,
        prompt: str = "Describe this image in detail.",
        model: str | None = None,
    ) -> str:
        """
        Возвращает текстовое описание изображения через LLaVA в Ollama.

        Позволяет не-vision LLM работать с изображениями: конвертер описывает
        картинку текстом, который затем передаётся в основную модель.

        Ollama должен быть запущен локально. Модель скачивается командой:
            ollama pull llava       # ~4 ГБ, универсальная
            ollama pull llava:13b   # ~8 ГБ, качественнее
            ollama pull llava-llama3  # альтернатива на базе Llama 3

        Параметры:
            path   -- абсолютный путь к изображению (png, jpg, jpeg, webp, gif, bmp).
            prompt -- текстовый запрос к модели (по умолчанию: общее описание).
                     Примеры: "What text is visible?", "List all objects.",
                              "Describe the chart data.", "What is the mood?"
            model  -- модель Ollama (None = использовать self.llava_model).

        Требует: запущенный Ollama + ollama pull llava
        """
        import urllib.request
        import urllib.error
        import json

        model = model or self.llava_model
        image_data = base64.b64encode(Path(path).read_bytes()).decode("utf-8")

        body: dict = {
            "model":      model,
            "prompt":     prompt,
            "images":     [image_data],
            "stream":     False,
            "keep_alive": _OLLAMA_KEEPALIVE,
        }
        if self.llava_cpu_only:
            # num_gpu=0 -- модель работает только на CPU, не вытесняя основную LLM из VRAM.
            # keep_alive держит модель загруженной между вызовами -- повторная инициализация
            # не нужна, пока не истечёт _OLLAMA_KEEPALIVE с момента последнего вызова.
            body["options"] = {"num_gpu": 0}

        payload = json.dumps(body).encode("utf-8")

        req = urllib.request.Request(
            f"{self.ollama_url}/api/generate",
            data=payload,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        try:
            with urllib.request.urlopen(req, timeout=120) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                return result.get("response", "").strip()
        except urllib.error.URLError as e:
            raise RuntimeError(_ERR_LLAVA_OLLAMA_UNAVAIL.format(model=model)) from e

    # -- Экспорт документов ---------------------------------------------------

    def export_pdf(
        self,
        content: str,
        output_path: str | None = None,
        title: str = "",
    ) -> str:
        """
        Конвертирует markdown или plain text в PDF через xhtml2pdf.

        Входной текст оборачивается в минимальный HTML с базовыми стилями
        (шрифт, отступы, переносы строк). Для сложного оформления --
        передавайте готовый HTML со своими стилями.

        Параметры:
            content     -- текст, markdown или HTML для конвертации.
            output_path -- путь к выходному PDF (None -> tempdir/export_<timestamp>.pdf).
            title       -- заголовок документа (добавляется как <title> и <h1>).

        Возвращает абсолютный путь к созданному PDF.
        Требует: pip install xhtml2pdf
        """
        if _pisa is None:
            raise RuntimeError(_ERR_EXPORT_PDF_NO_LIB)

        from datetime import datetime as _dt

        stripped = content.strip()
        if stripped.startswith("<") and ("</html>" in stripped.lower() or "<body" in stripped.lower()):
            html = content
        else:
            try:
                import markdown as _md
                body_html = _md.markdown(content, extensions=["tables", "fenced_code"])
            except ImportError:
                escaped = content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                body_html = "<p>" + escaped.replace("\n\n", "</p><p>").replace("\n", "<br>") + "</p>"

            h1 = f"<h1>{title}</h1>\n" if title else ""
            html = f"""<!DOCTYPE html>
<html><head>
<meta charset="utf-8">
<title>{title or 'Document'}</title>
<style>
  body {{ font-family: sans-serif; font-size: 12pt; line-height: 1.6;
         margin: 2cm; color: #222; }}
  h1, h2, h3 {{ color: #111; }}
  pre {{ background: #f5f5f5; padding: 0.8em;
        white-space: pre-wrap; font-size: 10pt; }}
  code {{ background: #f0f0f0; padding: 0.1em 0.3em; }}
  table {{ border-collapse: collapse; width: 100%; }}
  th, td {{ border: 1px solid #ccc; padding: 4px 8px; text-align: left; }}
  th {{ background: #eee; }}
</style>
</head><body>
{h1}{body_html}
</body></html>"""

        if output_path:
            out = Path(output_path)
        else:
            ts = _dt.now().strftime("%Y%m%d_%H%M%S_%f")
            out = Path(tempfile.gettempdir()) / f"export_{ts}.pdf"

        out.parent.mkdir(parents=True, exist_ok=True)

        with open(str(out), "wb") as f:
            result = _pisa.CreatePDF(html, dest=f, encoding="utf-8")

        if result.err:
            raise RuntimeError(f"export_pdf: xhtml2pdf conversion error (code {result.err})")

        return str(out)

    def export_docx(
        self,
        content: str,
        output_path: str | None = None,
        title: str = "",
    ) -> str:
        """
        Создаёт простой DOCX-документ из plain text через python-docx.

        Каждый непустой абзац входного текста становится параграфом документа.
        Строки, начинающиеся с # / ## / ###, распознаются как заголовки.
        Для сложного форматирования, таблиц, стилей -- используйте fastskills docx.

        Параметры:
            content     -- текст документа (поддерживаются # заголовки).
            output_path -- путь к выходному .docx (None -> tempdir/export_<timestamp>.docx).
            title       -- заголовок документа (добавляется как Heading 1 в начало).

        Возвращает абсолютный путь к созданному DOCX.
        Требует: pip install python-docx
        """
        if _docx is None:
            raise RuntimeError(_ERR_EXPORT_DOCX_NO_LIB)

        from datetime import datetime as _dt

        doc = _docx.Document()

        if title:
            doc.add_heading(title, level=1)

        for line in content.splitlines():
            stripped = line.rstrip()
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            else:
                doc.add_paragraph(stripped)

        if output_path:
            out = Path(output_path)
        else:
            ts = _dt.now().strftime("%Y%m%d_%H%M%S_%f")
            out = Path(tempfile.gettempdir()) / f"export_{ts}.docx"

        out.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(out))
        return str(out)

    def export_xlsx(
        self,
        rows: list,
        output_path: str | None = None,
        sheet_name: str = "Sheet1",
        headers: list | None = None,
    ) -> str:
        """
        Создаёт XLSX-таблицу из списка строк через openpyxl.

        rows -- список списков (каждый вложенный список = строка таблицы)
               или список словарей (ключи -> столбцы, headers формируются автоматически).
        Для сложного форматирования, формул, графиков -- используйте fastskills xlsx.

        Параметры:
            rows        -- данные: [[val, val, ...], ...] или [{key: val, ...}, ...].
            output_path -- путь к выходному .xlsx (None -> tempdir/export_<timestamp>.xlsx).
            sheet_name  -- название листа (по умолчанию "Sheet1").
            headers     -- заголовки столбцов (None -> без заголовков или из ключей словарей).

        Возвращает абсолютный путь к созданному XLSX.
        Требует: pip install openpyxl
        """
        if _openpyxl is None:
            raise RuntimeError(_ERR_EXPORT_XLSX_NO_LIB)

        from datetime import datetime as _dt

        wb = _openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name

        # Нормализуем rows: словари -> списки, заголовки из ключей если не заданы
        if rows and isinstance(rows[0], dict):
            if headers is None:
                headers = list(rows[0].keys())
            rows = [[row.get(h, "") for h in headers] for row in rows]

        if headers:
            ws.append(headers)

        for row in rows:
            ws.append([str(c) if c is not None else "" for c in row])

        if output_path:
            out = Path(output_path)
        else:
            ts = _dt.now().strftime("%Y%m%d_%H%M%S_%f")
            out = Path(tempfile.gettempdir()) / f"export_{ts}.xlsx"

        out.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(out))
        return str(out)

    # -- Интеграция с LangChain ------------------------------------------------

    def as_tools(self) -> list:
        """
        Возвращает список LangChain @tool-обёрток.
        Используется для встраивания конвертера напрямую в LangChain-агент
        без запуска отдельного MCP-процесса.

        Требует: pip install langchain
        """
        from langchain.tools import tool
        import json
        conv = self

        @tool
        def convert_pdf(path: str) -> str:
            """Извлекает текст из PDF-файла. path -- абсолютный путь."""
            return conv.pdf_to_text(path)

        @tool
        def convert_office(path: str) -> str:
            """Извлекает текст из Office-документа (.docx/.doc/.xlsx/.xls/.pptx). path -- абсолютный путь."""
            return conv.office_to_text(path)

        @tool
        def convert_html(path: str) -> str:
            """Извлекает читаемый текст из HTML-файла. path -- абсолютный путь."""
            return conv.html_to_text(path)

        @tool
        def transcribe_audio(path: str, language: str = "") -> str:
            """Транскрибирует аудио через Whisper. path -- абсолютный путь. language -- ISO 639-1 или пусто."""
            return conv.audio_to_transcript(path, language=language or None)

        @tool
        def transcribe_video(path: str) -> str:
            """Транскрибирует аудиодорожку видеофайла. path -- абсолютный путь."""
            result = conv.video_extract_audio_transcript(path)
            return result if result is not None else "[no audio track]"

        @tool
        def extract_video_frames(path: str, num_frames: int = 5) -> str:
            """Извлекает кадры из видео. path -- абсолютный путь. Возвращает JSON-массив base64-JPEG."""
            return json.dumps(conv.video_extract_frames_b64(path, num_frames=num_frames))

        @tool
        def describe_image(path: str, prompt: str = "Describe this image in detail.") -> str:
            """
            Описывает изображение текстом через LLaVA (Ollama). Используй для не-vision моделей.
            path -- абсолютный путь к изображению (png/jpg/jpeg/webp/gif/bmp).
            prompt -- что именно нужно описать (по умолчанию: общее описание).
            Требует запущенный Ollama и загруженную модель: ollama pull llava
            """
            return conv.describe_image(path, prompt=prompt)

        @tool
        def export_pdf(content: str, output_path: str = "", title: str = "") -> str:
            """
            Export plain text or markdown to a PDF file. Returns the path to the created file.
            Use for simple documents without custom styles or complex layout.
            For styled or structured PDFs use the fastskills pdf skill instead.

            Args:
                content:     Text, markdown, or HTML to export.
                output_path: Absolute path for the output PDF. Empty = auto-generate in temp dir.
                title:       Optional document title (added as H1 heading).
            Requires: pip install xhtml2pdf
            """
            return conv.export_pdf(content, output_path=output_path or None, title=title)

        @tool
        def export_docx(content: str, output_path: str = "", title: str = "") -> str:
            """
            Export plain text to a DOCX file. Lines starting with #/##/### become headings.
            Returns the path to the created file.
            Use for simple plain documents. For templates, styles, tables of contents,
            tracked changes, or complex formatting use the fastskills docx skill instead.

            Args:
                content:     Text to export. Lines with # / ## / ### become headings.
                output_path: Absolute path for the output .docx. Empty = auto-generate in temp dir.
                title:       Optional document title (added as Heading 1).
            Requires: pip install python-docx
            """
            return conv.export_docx(content, output_path=output_path or None, title=title)

        @tool
        def export_xlsx(
            rows: str,
            output_path: str = "",
            sheet_name: str = "Sheet1",
            headers: str = "",
        ) -> str:
            """
            Export tabular data to an XLSX file. Returns the path to the created file.
            Use for plain data tables. For formulas, charts, pivot tables, or conditional
            formatting use the fastskills xlsx skill instead.

            Args:
                rows:        JSON array of arrays [[v,v,...],[v,v,...]] or
                             array of objects [{col:val,...},{...}].
                output_path: Absolute path for the output .xlsx. Empty = auto-generate in temp dir.
                sheet_name:  Sheet name (default: Sheet1).
                headers:     JSON array of column names, e.g. ["Name","Age"].
                             Empty = no explicit headers (or auto-detected from dict keys).
            Requires: pip install openpyxl
            """
            import json as _json
            parsed_rows    = _json.loads(rows)
            parsed_headers = _json.loads(headers) if headers.strip() else None
            return conv.export_xlsx(
                parsed_rows,
                output_path=output_path or None,
                sheet_name=sheet_name or "Sheet1",
                headers=parsed_headers,
            )

        return [convert_pdf, convert_office, convert_html,
                transcribe_audio, transcribe_video, extract_video_frames,
                describe_image, export_pdf, export_docx, export_xlsx]


# =============================================================================
# MCP-сервер (модульный объект)
#
# `mcp` -- глобальный объект FastMCP, который обнаруживается автоматически
# при запуске через `fastmcp run file_converter.py` или `python file_converter.py`.
#
# Параметры Whisper берутся из переменных окружения при старте:
#   FILE_CONVERTER_CUDA=1          -> использовать GPU
#   FILE_CONVERTER_MODEL=large     -> размер модели (по умолчанию turbo)
#   FILE_CONVERTER_LANGUAGE=ru     -> язык по умолчанию (по умолчанию автодетект)
#   FILE_CONVERTER_LLAVA_MODEL=llava:13b -> модель LLaVA в Ollama (по умолчанию llava:7b)
#   FILE_CONVERTER_OLLAMA_URL=...  -> адрес Ollama API (по умолчанию http://localhost:11434)
# =============================================================================

try:
    from fastmcp import FastMCP as _FastMCP

    # Читаем конфигурацию из env -- удобно при запуске через Claude Desktop,
    # где передать аргументы командной строки можно только через "args" в JSON-конфиге,
    # а env-переменные задаются через секцию "env".
    _conv = FileConverter(
        use_cuda         = os.environ.get("FILE_CONVERTER_CUDA", "").lower() in ("1", "true", "yes"),
        whisper_model    = os.environ.get("FILE_CONVERTER_MODEL", "turbo"),
        whisper_language = os.environ.get("FILE_CONVERTER_LANGUAGE") or None,
        llava_model      = os.environ.get("FILE_CONVERTER_LLAVA_MODEL", DEFAULT_MODEL),
        ollama_url       = os.environ.get("FILE_CONVERTER_OLLAMA_URL", "http://localhost:11434"),
        llava_cpu_only   = os.environ.get("FILE_CONVERTER_LLAVA_CPU_ONLY", "1").lower() not in ("0", "false", "no"),
    )

    mcp = _FastMCP(
        "file-converter",
        instructions=(
            "Converts local files to text. "
            "All tools accept an absolute path to a file on the machine where the server runs. "
            "Supports: PDF, Office (docx/xlsx/pptx), HTML, audio (Whisper, local), video."
        ),
    )

    @mcp.tool()
    def convert_pdf(path: str) -> str:
        """
        Extract text from a PDF file using pypdf.

        Args:
            path: Absolute path to the PDF file.

        Returns text extracted page by page.
        Requires: pip install pypdf
        """
        return _conv.pdf_to_text(path)

    @mcp.tool()
    def convert_office(path: str) -> str:
        """
        Extract text from an Office document.

        Supported formats: .docx, .doc, .xlsx, .xls, .pptx
        Args:
            path: Absolute path to the file.

        For .doc files LibreOffice must be installed and available in PATH.
        Requires: pip install python-docx openpyxl python-pptx
        """
        return _conv.office_to_text(path)

    @mcp.tool()
    def convert_html(path: str) -> str:
        """
        Extract readable text from an HTML file (strips tags, scripts, styles).

        Args:
            path: Absolute path to the HTML file.

        Requires: pip install "beautifulsoup4[lxml]"
        """
        return _conv.html_to_text(path)

    @mcp.tool()
    def transcribe_audio(
        path: str,
        language: str = "",
        model_size: str = "",
        word_timestamps: bool = False,
    ) -> str:
        """
        Transcribe an audio file to text locally using faster-whisper (no API, no internet).

        The Whisper model is downloaded from HuggingFace on first use (~tens to hundreds of MB
        depending on model size) and cached locally for subsequent calls.

        Args:
            path:            Absolute path to the audio file (mp3, wav, ogg, m4a, flac, aac).
            language:        ISO 639-1 language code (ru, en, de, ...). Empty = auto-detect.
            model_size:      Whisper model size: tiny/base/small/medium/large/turbo.
                             Empty = use server default (turbo unless overridden via env).
            word_timestamps: If true, each word is prefixed with its timestamp: [0.00s -> 0.54s].

        Requires: pip install faster-whisper
        """
        return _conv.audio_to_transcript(
            path,
            model_size=model_size or None,
            language=language or None,
            word_timestamps=word_timestamps,
        )

    @mcp.tool()
    def transcribe_video(path: str, language: str = "") -> str:
        """
        Extract and transcribe the audio track of a video file using faster-whisper.

        Args:
            path:     Absolute path to the video file (mp4, mov, avi, webm).
            language: ISO 639-1 language code. Empty = auto-detect.

        Returns the transcript, or '[no audio track]' if the video has no audio.
        Requires: pip install faster-whisper moviepy
        """
        result = _conv.video_extract_audio_transcript(path)
        return result if result is not None else "[no audio track]"

    @mcp.tool()
    def extract_video_frames(path: str, num_frames: int = 5) -> str:
        """
        Extract evenly spaced frames from a video file.

        Args:
            path:       Absolute path to the video file.
            num_frames: Number of frames to extract (default: 5).

        Returns a JSON array of base64-encoded JPEG strings.
        Requires: pip install opencv-python
        """
        import json
        return json.dumps(_conv.video_extract_frames_b64(path, num_frames=num_frames))

    @mcp.tool()
    def describe_image(
        path: str,
        prompt: str = "Describe this image in detail.",
        model: str = "",
    ) -> str:
        """
        Generate a text description of an image using LLaVA via Ollama.

        Use this tool to enable non-vision LLMs to understand images: the image
        is described as text, which is then passed to the main language model.

        Args:
            path:   Absolute path to the image file (png, jpg, jpeg, webp, gif, bmp).
            prompt: What to ask about the image (default: general description).
                    Examples: "What text is visible?", "List all objects in the scene.",
                              "Describe the data shown in this chart.", "What is the mood?"
            model:  Ollama model to use. Empty = server default (llava unless overridden via env).
                    Other options: llava:13b, llava:34b, llava-llama3, bakllava.

        Requires: Ollama running locally + model pulled: ollama pull llava
        """
        return _conv.describe_image(path, prompt=prompt, model=model or None)

    @mcp.tool()
    def export_pdf(content: str, output_path: str = "", title: str = "") -> str:
        """
        Export plain text or markdown to a PDF file. Returns the path to the created file.
        Use for simple documents without custom styles or complex layout.
        For styled or structured PDFs use the fastskills pdf skill instead.

        Args:
            content:     Text, markdown, or HTML to export.
            output_path: Absolute path for the output PDF. Empty = auto-generate in temp dir.
            title:       Optional document title (added as H1 heading).

        Requires: pip install xhtml2pdf
        """
        return _conv.export_pdf(content, output_path=output_path or None, title=title)

    @mcp.tool()
    def export_docx(content: str, output_path: str = "", title: str = "") -> str:
        """
        Export plain text to a DOCX file. Lines starting with #/##/### become headings.
        Returns the path to the created file.
        Use for simple plain documents. For templates, styles, tables of contents,
        tracked changes, or complex formatting use the fastskills docx skill instead.

        Args:
            content:     Text to export. Lines with # / ## / ### become headings.
            output_path: Absolute path for the output .docx. Empty = auto-generate in temp dir.
            title:       Optional document title (added as Heading 1).

        Requires: pip install python-docx
        """
        return _conv.export_docx(content, output_path=output_path or None, title=title)

    @mcp.tool()
    def export_xlsx(
        rows: str,
        output_path: str = "",
        sheet_name: str = "Sheet1",
        headers: str = "",
    ) -> str:
        """
        Export tabular data to an XLSX file. Returns the path to the created file.
        Use for plain data tables. For formulas, charts, pivot tables, or conditional
        formatting use the fastskills xlsx skill instead.

        Args:
            rows:        JSON array of arrays [[v,v,...],[v,v,...]] or
                         array of objects [{col:val,...},{...}].
            output_path: Absolute path for the output .xlsx. Empty = auto-generate in temp dir.
            sheet_name:  Sheet name (default: Sheet1).
            headers:     JSON array of column names, e.g. ["Name","Age"].
                         Empty = no explicit headers (or auto-detected from dict keys).

        Requires: pip install openpyxl
        """
        import json as _json
        parsed_rows    = _json.loads(rows)
        parsed_headers = _json.loads(headers) if headers.strip() else None
        return _conv.export_xlsx(
            parsed_rows,
            output_path=output_path or None,
            sheet_name=sheet_name or "Sheet1",
            headers=parsed_headers,
        )

except ImportError:
    # fastmcp не установлен -- модульный объект mcp недоступен,
    # но FileConverter и as_tools() работают без него.
    mcp = None  # type: ignore


# =============================================================================
# Точка входа CLI
# =============================================================================

def main():
    """
    Запускает сервер в SSE-режиме.

    Используйте этот режим для подключения по HTTP из LangChain/AIList:
        await ai.mcp_connect("file-converter", url="http://127.0.0.1:8010/sse")

    Для Claude Desktop предпочтительнее stdio-режим через fastmcp:
        fastmcp run file_converter.py
    """
    import argparse

    parser = argparse.ArgumentParser(
        description="File Converter MCP server",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Environment variables (alternative to CLI args, useful for Claude Desktop):
  FILE_CONVERTER_CUDA=1        use GPU for Whisper
  FILE_CONVERTER_MODEL=large   Whisper model size
  FILE_CONVERTER_LANGUAGE=ru   default transcription language
  FILE_CONVERTER_LLAVA_MODEL=llava:7b  LLaVA model in Ollama
  FILE_CONVERTER_OLLAMA_URL=http://...  Ollama API URL
        """,
    )
    parser.add_argument("--host",  default="127.0.0.1", help="Bind address (default: 127.0.0.1)")
    parser.add_argument("--port",  type=int, default=8010, help="Port (default: 8010)")
    parser.add_argument("--cuda",  action="store_true", help="Use GPU for Whisper")
    parser.add_argument("--model", default="turbo", help="Whisper model size (default: turbo)")
    parser.add_argument("--lang",         default="",                         help="Default transcription language, e.g. ru (default: auto-detect)")
    parser.add_argument("--llava",        default="llava",                    help="LLaVA model in Ollama (default: llava)")
    parser.add_argument("--ollama-url",   default="http://localhost:11434",   help="Ollama API URL (default: http://localhost:11434)")
    parser.add_argument("--llava-use-gpu", action="store_true",               help="Allow GPU for vision model (default: CPU only, num_gpu=0)")
    args = parser.parse_args()

    if mcp is None:
        raise ImportError("fastmcp is not installed. Run: pip install fastmcp")

    _conv.use_cuda         = args.cuda
    _conv.whisper_model    = args.model
    _conv.whisper_language = args.lang or None
    _conv.llava_model      = args.llava
    _conv.ollama_url       = args.ollama_url.rstrip("/")
    _conv.llava_cpu_only   = not args.llava_use_gpu
    _conv._whisper_cache   = {}

    print(f"File Converter MCP server starting on {args.host}:{args.port}")
    print(f"Whisper: model={args.model}, cuda={args.cuda}, lang={args.lang or 'auto'}")
    print(f"LLaVA:   model={args.llava}, ollama={args.ollama_url}, cpu_only={_conv.llava_cpu_only}, keepalive={_OLLAMA_KEEPALIVE}")
    mcp.run(transport="sse", host=args.host, port=args.port)


if __name__ == "__main__":
    main()